"""Knowledge store — lexical retrieval over the documents/ folder + crawled site pages.

Runs fully in-process (stdlib only: no DB, no embeddings, no SDK) so it keeps working
even with AGENT_MAX_TOOL_ITERATIONS=0. Refresh is idempotent: every refresh rebuilds
the chunk index from disk and the site, so added/edited/removed sources are picked up
automatically. Content stays invisible to the UI — it is only used to ground answers.
"""
from __future__ import annotations

import csv
import html
import json
import math
import re
import threading
import time
import urllib.request
from collections import Counter
from collections.abc import Iterable
from html.parser import HTMLParser
from pathlib import Path
from typing import Any

from core.logging import get_logger
from core.settings import settings

from services.rag import chunk_text

logger = get_logger("services.knowledge")

USER_AGENT = "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AIKnowledgeBot/1.0"

_TERM_RE = re.compile(r"[a-z0-9]+")
_HTML_SKIP_TAGS = {"script", "style", "noscript", "header", "nav", "footer", "aside"}


def _tokenize(text: str) -> list[str]:
    return _TERM_RE.findall(text.lower())


class _TextExtractor(HTMLParser):
    """Collects visible text from HTML; drops script/style/noscript content."""

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self._parts: list[str] = []
        self._skip = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in _HTML_SKIP_TAGS:
            self._skip += 1

    def handle_endtag(self, tag: str) -> None:
        if tag in _HTML_SKIP_TAGS and self._skip:
            self._skip -= 1

    def handle_data(self, data: str) -> None:
        if not self._skip:
            self._parts.append(data)

    def text(self) -> str:
        return " ".join(piece.strip() for piece in self._parts if piece.strip())


def _flatten_json(obj: Any, prefix: str = "") -> Iterable[str]:
    if isinstance(obj, dict):
        for key, value in obj.items():
            label = f"{prefix}{key}: " if prefix else f"{key}: "
            yield from _flatten_json(value, label)
    elif isinstance(obj, list):
        for item in obj:
            yield from _flatten_json(item, prefix)
    else:
        yield f"{prefix}{obj}"


def _extract_text(path: Path) -> str | None:
    """Return plain text for a supported file, or None when unsupported/unreadable."""
    suffix = path.suffix.lower()
    try:
        if suffix in {".txt", ".md", ".markdown", ".rst", ".text", ".log"}:
            return path.read_text(encoding="utf-8", errors="replace")
        if suffix == ".csv":
            rows: list[str] = []
            with path.open(encoding="utf-8", errors="replace", newline="") as fh:
                for row in csv.reader(fh):
                    rows.append(" | ".join(cell.strip() for cell in row))
            return "\n".join(rows)
        if suffix == ".json":
            body = json.loads(path.read_text(encoding="utf-8", errors="replace"))
            return "\n".join(_flatten_json(body))
        if suffix in {".html", ".htm", ".xml"}:
            parser = _TextExtractor()
            parser.feed(path.read_text(encoding="utf-8", errors="replace"))
            return parser.text()
        if suffix == ".pdf":
            try:
                from pypdf import PdfReader  # type: ignore[import-not-found]
            except ImportError:
                logger.warning("knowledge_skip_pdf", reason="pypdf not installed", file=str(path))
                return None
            reader = PdfReader(str(path))
            return "\n".join((page.extract_text() or "") for page in reader.pages)
        if suffix == ".docx":
            try:
                from docx import Document  # type: ignore[import-not-found]
            except ImportError:
                logger.warning("knowledge_skip_docx", reason="python-docx not installed", file=str(path))
                return None
            return "\n".join(paragraph.text for paragraph in Document(str(path)).paragraphs)
        if suffix == ".xlsx":
            try:
                from openpyxl import load_workbook  # type: ignore[import-not-found]
            except ImportError:
                logger.warning("knowledge_skip_xlsx", reason="openpyxl not installed", file=str(path))
                return None
            wb = load_workbook(str(path), read_only=True, data_only=True)
            lines: list[str] = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    lines.append(" | ".join("" if cell is None else str(cell) for cell in row))
            wb.close()
            return "\n".join(lines)
        if suffix == ".pptx":
            try:
                from pptx import Presentation  # type: ignore[import-not-found]
            except ImportError:
                logger.warning("knowledge_skip_pptx", reason="python-pptx not installed", file=str(path))
                return None
            prs = Presentation(str(path))
            return "\n".join(
                shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
            )
    except Exception as exc:  # a corrupt file must never kill a refresh
        logger.warning("knowledge_extract_failed", file=str(path), error=str(exc))
        return None
    return None


def _http_get(url: str) -> str | None:
    request = urllib.request.Request(
        url, headers={"User-Agent": USER_AGENT, "Accept": "text/html,application/xhtml+xml"}
    )
    try:
        with urllib.request.urlopen(request, timeout=20) as resp:
            raw = resp.read(2 * 1024 * 1024)
    except Exception as exc:
        logger.warning("knowledge_fetch_failed", url=url, error=str(exc))
        return None
    return raw.decode("utf-8", errors="replace")


def _sitemap_urls(body: str) -> list[str]:
    return re.findall(r"<loc>\s*([^<]+?)\s*</loc>", body, flags=re.IGNORECASE)


def _page_title(body: str, fallback: str) -> str:
    match = re.search(r"<title[^>]*>(.*?)</title>", body, flags=re.IGNORECASE | re.DOTALL)
    return html.unescape(match.group(1)).strip() if match else fallback


class KnowledgeStore:
    """In-memory chunk index with BM25 retrieval. Rebuilt on every refresh."""

    def __init__(self, *, docs_dir: Path | None = None, sites: list[str] | None = None) -> None:
        self.docs_dir = Path(settings.knowledge_docs_dir) if docs_dir is None else docs_dir
        self.sites = list(settings.knowledge_sites) if sites is None else sites
        self._lock = threading.Lock()
        self._chunks: list[dict[str, Any]] = []
        self._sources: list[str] = []
        self._site_chunks: list[dict[str, Any]] = []
        self._stats: dict[str, Any] = {"last_refresh": None, "last_error": None, "duration_ms": None}

    # -- refresh ---------------------------------------------------------------

    def _finalize(
        self,
        chunks: list[dict[str, Any]],
        sources: list[str],
        *,
        started: float,
        include_site: bool,
    ) -> None:
        if include_site:
            with self._lock:
                site_chunks = list(self._site_chunks)
            chunks.extend(dict(item) for item in site_chunks)
            sources.extend(item["source"] for item in site_chunks)
        for index, chunk in enumerate(chunks):
            chunk["id"] = index
        with self._lock:
            self._chunks = chunks
            self._sources = sources
        self._stats.update(
            last_refresh=time.time(),
            last_error=None,
            duration_ms=int((time.monotonic() - started) * 1000),
            chunks=len(chunks),
        )
        logger.info(
            "knowledge_refreshed",
            chunks=len(chunks),
            sources=len(sources),
            duration_ms=self._stats["duration_ms"],
            site=include_site,
        )

    def refresh(self) -> None:
        started = time.monotonic()
        chunks: list[dict[str, Any]] = []
        sources: list[str] = []
        site_chunks: list[dict[str, Any]] = []
        try:
            chunks.extend(self._scan_documents(sources))
            site_chunks = self._crawl_site(sources)
            with self._lock:
                self._site_chunks = site_chunks
        except Exception as exc:
            logger.exception("knowledge_refresh_failed")
            self._stats["last_error"] = str(exc)
        self._finalize(chunks, sources, started=started, include_site=True)

    def refresh_documents_only(self) -> None:
        """Rebuild only the documents/ chunks, keeping cached site chunks.

        Used by the admin portal so uploads/edit/delete are reflected in
        retrieval without re-crawling the website."""
        started = time.monotonic()
        chunks: list[dict[str, Any]] = []
        sources: list[str] = []
        try:
            chunks.extend(self._scan_documents(sources))
        except Exception as exc:
            logger.exception("knowledge_refresh_docs_failed")
            self._stats["last_error"] = str(exc)
        self._finalize(chunks, sources, started=started, include_site=True)

    def _scan_documents(self, sources: list[str]) -> list[dict[str, Any]]:
        chunks: list[dict[str, Any]] = []
        if not self.docs_dir.is_dir():
            logger.warning("knowledge_docs_missing", path=str(self.docs_dir))
            return chunks
        for path in sorted(self.docs_dir.rglob("*")):
            if not path.is_file():
                continue
            text = _extract_text(path)
            if not text:
                continue
            rel = str(path.relative_to(self.docs_dir)).replace("\\", "/")
            source = f"documents/{rel}"
            sources.append(source)
            for piece in chunk_text(text):
                chunks.append({"source": source, "title": rel, "kind": "document", "text": piece})
        return chunks

    def _crawl_site(self, sources: list[str]) -> list[dict[str, Any]]:
        chunks: list[dict[str, Any]] = []
        for url in self._discover_urls():
            body = _http_get(url)
            if not body:
                continue
            text = _extract_html_text(body)
            if not text or len(text) < 200:
                continue
            title = _page_title(body, url)
            sources.append(url)
            for piece in chunk_text(text):
                chunks.append({"source": url, "title": title, "kind": "web", "text": piece})
        return chunks

    def _discover_urls(self) -> list[str]:
        seen: set[str] = set()
        for site in self.sites:
            base = site.rstrip("/")
            for hint in ("sitemap_index.xml", "sitemap.xml"):
                body = _http_get(f"{base}/{hint}")
                if not body:
                    continue
                locs = _sitemap_urls(body)
                pages = [loc for loc in locs if "sitemap" not in loc]
                for sub in [loc for loc in locs if "sitemap" in loc]:
                    sub_body = _http_get(sub)
                    if sub_body:
                        pages.extend(_sitemap_urls(sub_body))
                for page in pages:
                    if page.startswith(("http://", "https://")) and self._same_origin(base, page):
                        seen.add(page)
                if seen:
                    return sorted(seen)[: settings.knowledge_max_site_pages]
        for site in self.sites:
            seen.add(site.rstrip("/"))
        return sorted(seen)

    @staticmethod
    def _same_origin(base: str, page: str) -> bool:
        host = base.split("//")[-1].split("/")[0]
        return host in page

    # -- retrieval -------------------------------------------------------------

    def search(self, query: str, top_k: int = 5) -> list[dict[str, Any]]:
        with self._lock:
            chunks = list(self._chunks)
        tokens = _tokenize(query)
        if not chunks or not tokens:
            return []
        return self._bm25(chunks, tokens, top_k)

    @staticmethod
    def _bm25(chunks: list[dict[str, Any]], tokens: list[str], top_k: int) -> list[dict[str, Any]]:
        avgdl = sum(len(chunk["text"]) for chunk in chunks) / len(chunks)
        total = len(chunks)
        tf: dict[int, Counter[str]] = {}
        df: Counter[str] = Counter()
        for index, chunk in enumerate(chunks):
            counts = Counter(_tokenize(chunk["text"]))
            tf[index] = counts
            for term in counts:
                df[term] += 1
        k1, b = 1.2, 0.75
        scored: list[tuple[float, dict[str, Any]]] = []
        for index, chunk in enumerate(chunks):
            length = len(chunk["text"])
            score = 0.0
            for term in tokens:
                freq = tf[index].get(term, 0)
                if not freq:
                    continue
                idf = math.log(1 + (total - df[term] + 0.5) / (df[term] + 0.5))
                score += idf * (freq * (k1 + 1)) / (freq + k1 * (1 - b + b * length / avgdl))
            if score > 0:
                scored.append((score, chunk))
        scored.sort(key=lambda pair: pair[0], reverse=True)
        return [chunk for _, chunk in scored[:top_k]]

    # -- status ----------------------------------------------------------------

    def status(self) -> dict[str, Any]:
        with self._lock:
            sources = list(self._sources)
            stats = dict(self._stats)
        return {
            "docs_dir": str(self.docs_dir),
            "sites": list(self.sites),
            "sources": sources,
            "source_count": len(sources),
            "chunk_count": stats.get("chunks", 0),
            "last_refresh": stats.get("last_refresh"),
            "last_error": stats.get("last_error"),
            "duration_ms": stats.get("duration_ms"),
        }


def _extract_html_text(body: str) -> str:
    parser = _TextExtractor()
    parser.feed(body)
    return parser.text()
